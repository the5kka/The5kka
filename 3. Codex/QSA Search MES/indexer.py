import logging
import os
import traceback

import fitz
from docx import Document
from pptx import Presentation

import db


TARGET_EXTENSIONS = {".pptx", ".ppt", ".docx", ".doc", ".pdf"}
LOG_FILE_NAME = "index_errors.log"


def setup_logger():
    """읽기 실패 파일을 별도 로그로 남기기 위한 로거를 준비합니다."""
    log_path = os.path.join(db.get_app_dir(), LOG_FILE_NAME)
    logging.basicConfig(
        filename=log_path,
        level=logging.ERROR,
        format="%(asctime)s [%(levelname)s] %(message)s",
        encoding="utf-8",
    )


def normalize_text(text):
    """검색 미리보기가 깨지지 않도록 공백을 정리합니다."""
    if not text:
        return ""
    return " ".join(str(text).replace("\x00", " ").split())


def list_target_files(root_folder):
    """지정 폴더와 하위 폴더에서 대상 문서 파일을 찾습니다."""
    root_folder = db.normalize_path(root_folder)
    files = []
    for current_root, _, names in os.walk(root_folder):
        for name in names:
            if name.startswith("~$"):
                continue
            ext = os.path.splitext(name)[1].lower()
            if ext in TARGET_EXTENSIONS:
                files.append(db.normalize_path(os.path.join(current_root, name)))
    return files


def make_record(file_path, location, content, modified_time):
    """DB 저장용 문서 레코드를 생성합니다."""
    file_path = db.normalize_path(file_path)
    ext = os.path.splitext(file_path)[1].lower().replace(".", "").upper()
    return {
        "file_name": os.path.basename(file_path),
        "file_path": file_path,
        "file_type": ext,
        "location": location,
        "content": normalize_text(content),
        "modified_time": modified_time,
    }


def extract_pptx(file_path, modified_time):
    """PPTX 파일을 슬라이드 단위로 텍스트 추출합니다."""
    records = []
    prs = Presentation(file_path)
    for idx, slide in enumerate(prs.slides, start=1):
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        parts.append(cell.text)
        text = normalize_text("\n".join(parts))
        if text:
            records.append(make_record(file_path, f"Slide {idx}", text, modified_time))
    return records


def extract_docx(file_path, modified_time):
    """DOCX 파일을 문단 및 표 텍스트 기준으로 추출합니다."""
    records = []
    document = Document(file_path)

    for idx, paragraph in enumerate(document.paragraphs, start=1):
        text = normalize_text(paragraph.text)
        if text:
            records.append(make_record(file_path, f"Paragraph {idx}", text, modified_time))

    table_index = 1
    for table in document.tables:
        table_parts = []
        for row in table.rows:
            for cell in row.cells:
                table_parts.append(cell.text)
        text = normalize_text("\n".join(table_parts))
        if text:
            records.append(make_record(file_path, f"Table {table_index}", text, modified_time))
        table_index += 1

    return records


def extract_pdf(file_path, modified_time):
    """PDF 파일을 페이지 단위로 텍스트 추출합니다."""
    records = []
    pdf = fitz.open(file_path)
    try:
        for idx, page in enumerate(pdf, start=1):
            text = normalize_text(page.get_text("text"))
            if text:
                records.append(make_record(file_path, f"Page {idx}", text, modified_time))
    finally:
        pdf.close()
    return records


def extract_doc_with_word(file_path, modified_time):
    """오래된 DOC 파일을 MS Word COM으로 열어 텍스트를 추출합니다."""
    import pythoncom
    import win32com.client

    pythoncom.CoInitialize()
    word = None
    document = None
    records = []
    try:
        word = win32com.client.DispatchEx("Word.Application")
        word.Visible = False
        document = word.Documents.Open(file_path, ReadOnly=True, AddToRecentFiles=False)
        for idx, paragraph in enumerate(document.Paragraphs, start=1):
            text = normalize_text(paragraph.Range.Text)
            if text:
                records.append(make_record(file_path, f"Paragraph {idx}", text, modified_time))
    finally:
        try:
            if document is not None:
                document.Close(False)
        except Exception:
            pass
        try:
            if word is not None:
                word.Quit()
        except Exception:
            pass
        pythoncom.CoUninitialize()
    return records


def extract_ppt_with_powerpoint(file_path, modified_time):
    """오래된 PPT 파일을 MS PowerPoint COM으로 열어 텍스트를 추출합니다."""
    import pythoncom
    import win32com.client

    pythoncom.CoInitialize()
    powerpoint = None
    presentation = None
    records = []
    try:
        powerpoint = win32com.client.DispatchEx("PowerPoint.Application")
        presentation = powerpoint.Presentations.Open(
            file_path,
            WithWindow=False,
            ReadOnly=True,
            Untitled=False,
        )
        for slide_idx in range(1, presentation.Slides.Count + 1):
            slide = presentation.Slides(slide_idx)
            parts = []
            for shape_idx in range(1, slide.Shapes.Count + 1):
                shape = slide.Shapes(shape_idx)
                try:
                    if shape.HasTextFrame and shape.TextFrame.HasText:
                        parts.append(shape.TextFrame.TextRange.Text)
                except Exception:
                    continue
            text = normalize_text("\n".join(parts))
            if text:
                records.append(make_record(file_path, f"Slide {slide_idx}", text, modified_time))
    finally:
        try:
            if presentation is not None:
                presentation.Close()
        except Exception:
            pass
        try:
            if powerpoint is not None:
                powerpoint.Quit()
        except Exception:
            pass
        pythoncom.CoUninitialize()
    return records


def extract_file(file_path):
    """확장자에 맞는 방식으로 문서 내부 텍스트를 추출합니다."""
    file_path = db.normalize_path(file_path)
    modified_time = os.path.getmtime(file_path)
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pptx":
        return extract_pptx(file_path, modified_time)
    if ext == ".ppt":
        return extract_ppt_with_powerpoint(file_path, modified_time)
    if ext == ".docx":
        return extract_docx(file_path, modified_time)
    if ext == ".doc":
        return extract_doc_with_word(file_path, modified_time)
    if ext == ".pdf":
        return extract_pdf(file_path, modified_time)
    return []


def should_skip(file_path):
    """수정일이 동일한 파일은 재처리하지 않도록 판단합니다."""
    file_path = db.normalize_path(file_path)
    current_modified = os.path.getmtime(file_path)
    indexed_modified = db.get_indexed_modified_time(file_path)
    return indexed_modified is not None and abs(float(indexed_modified) - float(current_modified)) < 0.001


def build_index(root_folder, progress_callback=None, status_callback=None, stop_event=None):
    """문서 인덱스를 생성합니다. GUI 멈춤 방지를 위해 별도 스레드에서 호출됩니다."""
    root_folder = db.normalize_path(root_folder)
    setup_logger()
    db.init_db()
    deleted_count = db.remove_deleted_files()

    files = list_target_files(root_folder)
    total = len(files)
    indexed_count = 0
    skipped_count = 0
    failed_count = 0

    for index, file_path in enumerate(files, start=1):
        if stop_event is not None and stop_event.is_set():
            break

        file_name = os.path.basename(file_path)
        if status_callback:
            status_callback(f"처리 중: {file_name}")
        if progress_callback:
            progress_callback(index, total, file_name)

        try:
            if should_skip(file_path):
                skipped_count += 1
                continue

            records = extract_file(file_path)
            conn = db.connect_db()
            try:
                db.delete_file_records(file_path, conn=conn)
                conn.commit()
            finally:
                conn.close()
            db.insert_records(records)
            indexed_count += 1
        except Exception as exc:
            failed_count += 1
            logging.error("파일 읽기 실패: %s\n%s", file_path, traceback.format_exc())
            if status_callback:
                status_callback(f"읽기 실패: {file_name} ({exc})")

    if progress_callback:
        progress_callback(total, total, "")
    return {
        "total": total,
        "indexed": indexed_count,
        "skipped": skipped_count,
        "failed": failed_count,
        "deleted": deleted_count,
    }
